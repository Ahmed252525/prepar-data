import os
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import xml.etree.ElementTree as ET
from PIL import Image
import io
import zipfile
import re
import matplotlib.pyplot as plt
import numpy as np
from datetime import datetime
import json
import uuid

class PowerPointToMarkdownConverter:
    def __init__(self, pptx_file_path, output_dir="converted_presentation"):
        # Initialize converter with file path and output directory
        self.pptx_file_path = pptx_file_path
        self.output_dir = output_dir
        self.images_dir = os.path.join(output_dir, "images")
        self.charts_dir = os.path.join(output_dir, "charts")
        self.data_dir = os.path.join(output_dir, "data")
        self.doc_id = str(uuid.uuid4())
        
        # Create output directories
        for directory in [self.output_dir, self.images_dir, self.charts_dir, self.data_dir]:
            os.makedirs(directory, exist_ok=True)
        
        self.presentation = None
        self.chart_counter = 0
        self.image_counter = 0
        self.chart_metadata = []
    
    def convert(self):
        # Convert PowerPoint to Markdown
        try:
            if not os.path.exists(self.pptx_file_path):
                raise FileNotFoundError(f"No such file: {self.pptx_file_path}")
            self.presentation = Presentation(self.pptx_file_path)
            print(f"Converting presentation with {len(self.presentation.slides)} slides...")
            
            markdown_content = self._generate_markdown()
            
            markdown_file = os.path.join(self.output_dir, "presentation.md")
            with open(markdown_file, 'w', encoding='utf-8') as f:
                f.write(markdown_content)
            
            self._generate_chart_summary()
            self._generate_metadata()
            
            print(f"✅ Conversion complete! Files saved to: {self.output_dir}")
            print(f"📄 Main file: {markdown_file}")
            
            return markdown_file
            
        except Exception as e:
            print(f"Error during conversion: {str(e)}")
            return None
    
    def _generate_markdown(self):
        # Generate Markdown content for the presentation
        markdown_lines = []
        presentation_title = self._extract_presentation_title()
        markdown_lines.append(f"# {presentation_title}\n")
        markdown_lines.append(f"*Converted from PowerPoint on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*\n")
        markdown_lines.append(f"*Document ID: {self.doc_id}*\n")
        markdown_lines.append("---\n")
        
        markdown_lines.append("## Table of Contents\n")
        for i, slide in enumerate(self.presentation.slides, 1):
            slide_title = self._extract_slide_title(slide, i)
            markdown_lines.append(f"{i}. [{slide_title}](#slide-{i})\n")
        markdown_lines.append("\n---\n")
        
        for slide_num, slide in enumerate(self.presentation.slides, 1):
            slide_markdown = self._process_slide(slide, slide_num)
            markdown_lines.extend(slide_markdown)
            markdown_lines.append("\n---\n")
        
        return "\n".join(markdown_lines)
    
    def _extract_presentation_title(self):
        # Extract presentation title from first slide or file name
        if self.presentation.slides:
            first_slide = self.presentation.slides[0]
            title = self._extract_slide_title(first_slide, 1)
            if title != "Slide 1":
                return title
        return os.path.splitext(os.path.basename(self.pptx_file_path))[0]
    
    def _extract_slide_title(self, slide, slide_num):
        # Extract slide title from shapes or use default
        try:
            if slide.shapes.title and slide.shapes.title.text.strip():
                return slide.shapes.title.text.strip()
        except:
            pass
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                text = shape.text_frame.text.strip()
                if text and len(text) < 100:
                    return text
        return f"Slide {slide_num}"
    
    def _process_slide(self, slide, slide_num):
        # Process slide content
        markdown_lines = []
        slide_title = self._extract_slide_title(slide, slide_num)
        markdown_lines.append(f"## Slide {slide_num}: {slide_title} {{#slide-{slide_num}}}\n")
        
        slide_layout = getattr(slide.slide_layout, 'name', 'Unknown Layout')
        markdown_lines.append(f"*Layout: {slide_layout}*\n")
        
        if hasattr(slide, 'notes_slide') and slide.notes_slide:
            notes_text = self._extract_notes(slide.notes_slide)
            if notes_text:
                markdown_lines.append(f"**Presenter Notes:** {notes_text}\n")
        
        text_content = []
        media_content = []
        
        for shape_idx, shape in enumerate(slide.shapes):
            shape_content = self._process_shape(shape, slide_num, shape_idx)
            if shape_content:
                if shape_content['type'] in ['text', 'title']:
                    text_content.append(shape_content)
                else:
                    media_content.append(shape_content)
        
        for content in text_content:
            if content['content']:
                markdown_lines.extend(content['content'])
        
        if media_content:
            markdown_lines.append("\n### Charts and Media\n")
            for content in media_content:
                if content['content']:
                    markdown_lines.extend(content['content'])
        
        return markdown_lines
    
    def _process_shape(self, shape, slide_num, shape_idx):
        # Process individual shape based on type
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                return self._process_chart(shape, slide_num, shape_idx)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return self._process_image(shape, slide_num, shape_idx)
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                return self._process_table(shape, slide_num, shape_idx)
            elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or hasattr(shape, 'text_frame'):
                return self._process_text(shape, slide_num, shape_idx)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                return self._process_group(shape, slide_num, shape_idx)
            elif shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
                return self._process_ole_object(shape, slide_num, shape_idx)
            else:
                return self._process_generic_shape(shape, slide_num, shape_idx)
                
        except Exception as e:
            print(f"Error processing shape in slide {slide_num}: {str(e)}")
            return None
    
    def _process_chart(self, chart_shape, slide_num, shape_idx):
        # Process chart shape
        try:
            self.chart_counter += 1
            chart_filename = f"chart_{self.doc_id}_{slide_num}_{self.chart_counter}"
            
            chart_data = self._extract_chart_data(chart_shape)
            
            is_hybrid = self._is_table_bar_hybrid(chart_shape)
            if is_hybrid:
                chart_data['is_hybrid'] = True
                chart_data['table_labels'] = self._extract_table_labels(chart_shape)
            
            chart_image_path = self._create_chart_image(chart_data, chart_filename)
            
            data_file = os.path.join(self.data_dir, f"{chart_filename}.json")
            with open(data_file, 'w', encoding='utf-8') as f:
                json.dump(chart_data, f, indent=2, ensure_ascii=False)
            
            self.chart_metadata.append({
                'doc_id': self.doc_id,
                'slide_num': slide_num,
                'chart_id': self.chart_counter,
                'filename': chart_filename,
                'type': 'chart',
                'path': chart_image_path if chart_image_path else 'No chart image generated'
            })
            
            content = []
            content.append(f"#### Chart: {chart_data.get('title', 'Untitled Chart')}\n")
            
            if chart_image_path:
                rel_image_path = os.path.relpath(chart_image_path, self.output_dir)
                content.append(f"![{chart_data.get('title', 'Chart')}]({rel_image_path})\n")
            else:
                content.append("**Note:** Chart image could not be generated.\n")
            
            content.append(f"**Chart Type:** {chart_data.get('chart_type', 'Unknown')}\n")
            
            if chart_data.get('categories'):
                content.append(f"**Categories:** {', '.join(map(str, chart_data['categories'][:5]))}")
                if len(chart_data['categories']) > 5:
                    content.append("...")
                content.append("\n")
            
            if is_hybrid and chart_data.get('table_labels'):
                content.append(f"**Table Labels:** {', '.join(chart_data['table_labels'][:5])}\n")
            
            if chart_data.get('series_data'):
                content.append("**Data Series:**\n")
                for series in chart_data['series_data'][:3]:
                    series_name = series.get('name', 'Unnamed Series')
                    values_count = len(series.get('values', []))
                    content.append(f"- {series_name} ({values_count} values)\n")
                
                if len(chart_data['series_data']) > 3:
                    content.append(f"- ... and {len(chart_data['series_data']) - 3} more series\n")
            
            rel_data_path = os.path.relpath(data_file, self.output_dir)
            content.append(f"📊 [View Chart Data]({rel_data_path})\n")
            
            return {
                'type': 'chart',
                'content': content
            }
            
        except Exception as e:
            print(f"Error processing chart in slide {slide_num}: {str(e)}")
            return None
    
    def _is_table_bar_hybrid(self, chart_shape):
        # Check if chart is a table-bar hybrid
        try:
            chart_xml = chart_shape.chart.part.related_parts.get('/ppt/charts/chart1.xml')
            if chart_xml:
                xml_str = chart_xml.blob.decode('utf-8')
                return 'table' in xml_str.lower() and 'bar' in xml_str.lower()
            return False
        except:
            return False
    
    def _extract_table_labels(self, chart_shape):
        # Extract labels from hybrid chart
        try:
            labels = []
            chart = chart_shape.chart
            if hasattr(chart, 'plots') and chart.plots:
                for plot in chart.plots:
                    if hasattr(plot, 'data_labels') and plot.data_labels:
                        for label in plot.data_labels:
                            if hasattr(label, 'text_frame') and label.text_frame.text:
                                labels.append(label.text_frame.text.strip())
            return labels
        except:
            return []
    
    def _extract_chart_data(self, chart_shape):
        # Extract data from chart
        try:
            chart = chart_shape.chart
            chart_data = {
                'title': 'Untitled Chart',
                'chart_type': str(chart.chart_type) if hasattr(chart, 'chart_type') else 'Unknown',
                'categories': [],
                'series_data': [],
                'has_3d': False,
                'doc_id': self.doc_id
            }
            
            if chart.has_title and chart.chart_title:
                chart_data['title'] = chart.chart_title.text_frame.text.strip() if chart.chart_title.text_frame.text else 'Untitled Chart'
            
            # Handle 3D chart types
            chart_type_str = str(chart.chart_type).lower()
            chart_data['has_3d'] = any(keyword in chart_type_str for keyword in ['3d', 'three', 'dimensional'])
            
            # Map unsupported 3D chart types to supported types
            if 'pie3dchart' in chart_type_str:
                chart_data['chart_type'] = 'pie'
                chart_data['has_3d'] = True
            elif 'bar3dchart' in chart_type_str:
                chart_data['chart_type'] = 'bar'
                chart_data['has_3d'] = True
            
            if hasattr(chart, 'plots') and chart.plots:
                for plot in chart.plots:
                    if hasattr(plot, 'categories') and plot.categories:
                        categories = [str(cat) for cat in plot.categories if cat is not None]
                        chart_data['categories'] = categories
                    
                    if hasattr(plot, 'series'):
                        for series in plot.series:
                            series_info = {
                                'name': getattr(series, 'name', f'Series {len(chart_data["series_data"]) + 1}'),
                                'values': []
                            }
                            if hasattr(series, 'values') and series.values is not None:
                                series_info['values'] = [float(v) if v is not None else 0 for v in series.values]
                            chart_data['series_data'].append(series_info)
            
            return chart_data
            
        except Exception as e:
            print(f"Error extracting chart data: {str(e)}")
            return {'title': 'Chart', 'chart_type': 'Unknown', 'categories': [], 'series_data': [], 'has_3d': False, 'doc_id': self.doc_id}
    
    def _create_chart_image(self, chart_data, filename):
        # Create chart image
        try:
            if not isinstance(chart_data, dict):
                raise ValueError("chart_data must be a dictionary")
            
            plt.style.use('ggplot')
            fig = plt.figure(figsize=(10, 6))
            
            chart_type = chart_data.get('chart_type', '').lower()
            is_3d = chart_data.get('has_3d', False)
            
            if is_3d and 'pie' not in chart_type:
                ax = fig.add_subplot(111, projection='3d')
            else:
                ax = fig.add_subplot(111)
            
            categories = chart_data.get('categories', [])
            series_data = chart_data.get('series_data', [])  # Fixed typo
            
            if not series_data:
                ax.text(0.5, 0.5, 0.5 if is_3d else 0, 
                        f"Title: {chart_data.get('title', 'No Data Available')}", 
                        ha='center', va='center', fontsize=12, transform=ax.transAxes)
                ax.set_xlim(0, 1)
                ax.set_ylim(0, 1)
                if is_3d:
                    ax.set_zlim(0, 1)
            else:
                if 'bar' in chart_type or 'column' in chart_type:
                    self._create_bar_chart(ax, categories, series_data, is_3d)
                elif 'line' in chart_type:
                    self._create_line_chart(ax, categories, series_data)
                elif 'pie' in chart_type:
                    self._create_pie_chart(ax, categories, series_data, is_3d)
                elif 'scatter' in chart_type:
                    self._create_scatter_chart(ax, series_data, is_3d)
                else:
                    self._create_bar_chart(ax, categories, series_data, is_3d)
            
            ax.set_title(chart_data.get('title', 'Chart'), fontsize=14, fontweight='bold', pad=20)
            
            if is_3d:
                ax.text2D(0.02, 0.98, '3D Chart', transform=ax.transAxes, 
                         bbox=dict(boxstyle="round,pad=0.3", facecolor='white', edgecolor='black'),
                         verticalalignment='top', fontsize=10)
            
            plt.tight_layout()
            chart_path = os.path.join(self.charts_dir, f"{filename}.png")
            plt.savefig(chart_path, dpi=300, bbox_inches='tight', facecolor='white')
            plt.close(fig)
            
            return chart_path
            
        except Exception as e:
            print(f"Error creating chart image: {str(e)}")
            return None
    
    def _create_bar_chart(self, ax, categories, series_data, is_3d=False):
        # Create bar chart
        if not categories and series_data:
            categories = [f"Item {i+1}" for i in range(len(series_data[0].get('values', [])))]
        
        x = np.arange(len(categories))
        width = 0.8 / len(series_data) if len(series_data) > 1 else 0.6
        colors = plt.cm.Set3(np.linspace(0, 1, len(series_data)))
        
        has_label = False
        for i, series in enumerate(series_data):
            values = series.get('values', [])
            if values:
                offset = (i - len(series_data)/2 + 0.5) * width
                if is_3d:
                    z = np.zeros(len(values))
                    dx = np.ones(len(values)) * width
                    dy = np.ones(len(values))
                    bars = ax.bar3d(x + offset, z, z, dx, dy, values, color=colors[i], alpha=0.8)
                    for bar in bars:
                        bar.set_edgecolor('black')
                        bar.set_linewidth(0.5)
                else:
                    bars = ax.bar(x + offset, values, width, label=series.get('name', f'Series {i+1}'), 
                                  color=colors[i], alpha=0.8)
                    has_label = True
                    for bar in bars:
                        bar.set_edgecolor('black')
                        bar.set_linewidth(0.5)
        
        ax.set_xlabel('Categories')
        ax.set_ylabel('Values')
        if is_3d:
            ax.set_zlabel('Values')
        ax.set_xticks(x)
        ax.set_xticklabels(categories, rotation=45, ha='right')
        if has_label and len(series_data) > 1:
            ax.legend()
    
    def _create_line_chart(self, ax, categories, series_data):
        # Create line chart
        if not categories:
            categories = [f"Point {i+1}" for i in range(len(series_data[0].get('values', [])))]
        
        x = np.arange(len(categories))
        colors = plt.cm.Set1(np.linspace(0, 1, len(series_data)))
        
        has_label = False
        for i, series in enumerate(series_data):
            values = series.get('values', [])
            if values:
                ax.plot(x, values, marker='o', label=series.get('name', f'Series {i+1}'), 
                       color=colors[i], linewidth=2, markersize=4)
                has_label = True
        
        ax.set_xlabel('Categories')
        ax.set_ylabel('Values')
        ax.set_xticks(x)
        ax.set_xticklabels(categories, rotation=45, ha='right')
        ax.grid(True, alpha=0.3)
        if has_label and len(series_data) > 1:
            ax.legend()
    
    def _create_pie_chart(self, ax, categories, series_data, is_3d=False):
        # Create pie chart
        if series_data and series_data[0].get('values'):
            values = series_data[0]['values']
            labels = categories if categories else [f"Slice {i+1}" for i in range(len(values))]
            
            filtered_data = [(label, value) for label, value in zip(labels, values) if value > 0]
            if filtered_data:
                labels, values = zip(*filtered_data)
                colors = plt.cm.Set3(np.linspace(0, 1, len(values)))
                
                if is_3d:
                    wedges, texts, autotexts = ax.pie(values, labels=labels, autopct='%1.1f%%', 
                                                     colors=colors, startangle=90, 
                                                     shadow=True, explode=[0.1] * len(values))
                else:
                    wedges, texts, autotexts = ax.pie(values, labels=labels, autopct='%1.1f%%', 
                                                     colors=colors, startangle=90)
                
                for autotext in autotexts:
                    autotext.set_color('white')
                    autotext.set_fontweight('bold')
        
        ax.axis('equal')
    
    def _create_scatter_chart(self, ax, series_data, is_3d=False):
        # Create scatter chart
        colors = plt.cm.Set1(np.linspace(0, 1, len(series_data)))
        
        for i, series in enumerate(series_data[:2]):
            values = series.get('values', [])
            if values:
                if i == 0:
                    x_values = values
                    x_label = series.get('name', 'X Values')
                else:
                    y_values = values[:len(x_values)]
                    y_label = series.get('name', 'Y Values')
                    if is_3d:
                        z_values = [0] * len(x_values)
                        ax.scatter(x_values, y_values, z_values, color=colors[0], alpha=0.6, s=50)
                        ax.set_zlabel('Z')
                    else:
                        ax.scatter(x_values, y_values, color=colors[0], alpha=0.6, s=50)
                    ax.set_xlabel(x_label)
                    ax.set_ylabel(y_label)
                    break
        
        ax.grid(True, alpha=0.3)
    
    def _process_image(self, image_shape, slide_num, shape_idx):
        # Process image shape
        try:
            self.image_counter += 1
            image_filename = f"image_{self.doc_id}_{slide_num}_{self.image_counter}"
            image = image_shape.image
            image_ext = getattr(image, 'ext', 'png').lstrip('.')
            
            image_path = os.path.join(self.images_dir, f"{image_filename}.{image_ext}")
            with open(image_path, 'wb') as f:
                f.write(image.blob)
            
            self.chart_metadata.append({
                'doc_id': self.doc_id,
                'slide_num': slide_num,
                'image_id': self.image_counter,
                'filename': image_filename,
                'type': 'image',
                'path': image_path
            })
            
            rel_image_path = os.path.relpath(image_path, self.output_dir)
            content = [f"![Image]({rel_image_path})\n"]
            
            return {
                'type': 'image',
                'content': content
            }
            
        except Exception as e:
            print(f"Error processing image: {str(e)}")
            return None
    
    def _process_table(self, table_shape, slide_num, shape_idx):
        # Process table shape
        try:
            table = table_shape.table
            table_filename = f"table_{self.doc_id}_{slide_num}_{shape_idx}"
            
            table_data = {
                'doc_id': self.doc_id,
                'slide_num': slide_num,
                'rows': []
            }
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    cell_text = cell.text.strip().replace('\n', ' ').replace('|', '\\|')
                    row_data.append(cell_text)
                table_data['rows'].append(row_data)
            
            if not table_data['rows']:
                return None
            
            data_file = os.path.join(self.data_dir, f"{table_filename}.json")
            with open(data_file, 'w', encoding='utf-8') as f:
                json.dump(table_data, f, indent=2, ensure_ascii=False)
            
            self.chart_metadata.append({
                'doc_id': self.doc_id,
                'slide_num': slide_num,
                'table_id': shape_idx,
                'filename': table_filename,
                'type': 'table',
                'path': data_file
            })
            
            markdown_table = []
            markdown_table.append("| " + " | ".join(table_data['rows'][0]) + " |")
            markdown_table.append("|" + "---|" * len(table_data['rows'][0]))
            for row in table_data['rows'][1:]:
                markdown_table.append("| " + " | ".join(row) + " |")
            markdown_table.append("")
            
            rel_data_path = os.path.relpath(data_file, self.output_dir)
            markdown_table.append(f"📋 [View Table Data]({rel_data_path})\n")
            
            return {
                'type': 'table',
                'content': markdown_table
            }
            
        except Exception as e:
            print(f"Error processing table: {str(e)}")
            return None
    
    def _process_text(self, text_shape, slide_num, shape_idx):
        # Process text shape
        try:
            if not hasattr(text_shape, 'text_frame') or not text_shape.text_frame:
                return None
            
            text_content = []
            is_title = False
            try:
                if hasattr(text_shape, 'placeholder_format') and text_shape.placeholder_format:
                    is_title = 'title' in str(text_shape.placeholder_format.type).lower()
            except:
                pass  # Skip shapes without valid placeholder_format
            
            for paragraph in text_shape.text_frame.paragraphs:
                if paragraph.text.strip():
                    level = getattr(paragraph, 'level', 0)
                    text = paragraph.text.strip()
                    if is_title and level == 0:
                        text_content.append(f"### {text}\n")
                    elif level == 0:
                        text_content.append(f"{text}\n")
                    else:
                        bullet = "  " * level + "- "
                        text_content.append(f"{bullet}{text}\n")
            
            if text_content:
                text_content.append("")
            
            return {
                'type': 'title' if is_title else 'text',
                'content': text_content
            }
            
        except Exception as e:
            print(f"Error processing text: {str(e)}")
            return None
    
    def _process_group(self, group_shape, slide_num, shape_idx):
        # Process group shape
        try:
            group_content = []
            for i, shape in enumerate(group_shape.shapes):
                shape_content = self._process_shape(shape, slide_num, f"{shape_idx}_g{i}")
                if shape_content and shape_content.get('content'):
                    group_content.extend(shape_content['content'])
            
            if group_content:
                return {
                    'type': 'group',
                    'content': group_content
                }
            return None
            
        except Exception as e:
            print(f"Error processing group: {str(e)}")
            return None
    
    def _process_ole_object(self, ole_shape, slide_num, shape_idx):
        # Process embedded OLE object
        try:
            prog_id = getattr(ole_shape.ole_format, 'prog_id', 'Unknown')
            content = [f"**Embedded Object:** {prog_id}\n"]
            
            if hasattr(ole_shape.ole_format, 'blob'):
                ext = 'xlsx' if 'excel' in prog_id.lower() else 'docx' if 'word' in prog_id.lower() else 'bin'
                obj_filename = f"embedded_{self.doc_id}_{slide_num}_{shape_idx}.{ext}"
                obj_path = os.path.join(self.data_dir, obj_filename)
                
                with open(obj_path, 'wb') as f:
                    f.write(ole_shape.ole_format.blob)
                
                rel_obj_path = os.path.relpath(obj_path, self.output_dir)
                content.append(f"📎 [Download Embedded Object]({rel_obj_path})\n")
            
            return {
                'type': 'ole_object',
                'content': content
            }
            
        except Exception as e:
            print(f"Error processing OLE object: {str(e)}")
            return None
    
    def _process_generic_shape(self, shape, slide_num, shape_idx):
        # Process generic shape
        try:
            content = [f"*Shape: {shape.shape_type}*\n"]
            if hasattr(shape, 'text_frame') and shape.text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    content.append(f"{text}\n")
            
            return {
                'type': 'generic_shape',
                'content': content
            }
            
        except Exception as e:
            print(f"Error processing generic shape: {str(e)}")
            return None
    
    def _extract_notes(self, notes_slide):
        # Extract presenter notes
        try:
            if hasattr(notes_slide, 'notes_text_frame') and notes_slide.notes_text_frame:
                return notes_slide.notes_text_frame.text.strip()
        except:
            pass
        return ""
    
    def _generate_chart_summary(self):
        # Generate chart and table summary
        summary_file = os.path.join(self.output_dir, "charts_summary.md")
        
        with open(summary_file, 'w', encoding='utf-8') as f:
            f.write("# Charts and Tables Summary\n\n")
            f.write(f"Document ID: {self.doc_id}\n")
            f.write(f"Total charts extracted: {self.chart_counter}\n")
            f.write(f"Total images extracted: {self.image_counter}\n\n")
            
            f.write("## Chart and Table Metadata\n\n")
            for meta in sorted(self.chart_metadata, key=lambda x: (x['slide_num'], x.get('chart_id', x.get('image_id', x.get('table_id', 0))))):
                f.write(f"- Slide {meta['slide_num']}: {meta['type'].capitalize()} ({meta['filename']})\n")
                f.write(f"  - Path: [{meta['filename']}]({os.path.relpath(meta['path'], self.output_dir)})\n")
                f.write(f"  - Document ID: {meta['doc_id']}\n")
    
    def _generate_metadata(self):
        # Generate metadata JSON
        metadata = {
            'doc_id': self.doc_id,
            'source_file': os.path.basename(self.pptx_file_path),
            'conversion_date': datetime.now().isoformat(),
            'total_slides': len(self.presentation.slides) if self.presentation else 0,
            'charts_extracted': self.chart_counter,
            'images_extracted': self.image_counter,
            'chart_metadata': self.chart_metadata,
            'output_structure': {
                'main_file': 'presentation.md',
                'charts_directory': 'charts/',
                'images_directory': 'images/',
                'data_directory': 'data/',
                'summary_file': 'charts_summary.md'
            }
        }
        
        metadata_file = os.path.join(self.output_dir, "metadata.json")
        with open(metadata_file, 'w', encoding='utf-8') as f:
            json.dump(metadata, f, indent=2, ensure_ascii=False)

def convert_powerpoint_to_markdown(pptx_file, output_dir=None):
    # Main function to convert PowerPoint to Markdown
    if output_dir is None:
        base_name = os.path.splitext(os.path.basename(pptx_file))[0]
        output_dir = f"{base_name}_markdown"
    
    converter = PowerPointToMarkdownConverter(pptx_file, output_dir)
    return converter.convert()

if __name__ == "__main__":
    # Run conversion with fixed file name
    pptx_file = "test.pptx"
    if os.path.exists(pptx_file):
        print("🚀 Starting PowerPoint to Markdown conversion...")
        result = convert_powerpoint_to_markdown(pptx_file)
        if result:
            print(f"✅ Conversion successful!")
            print(f"📁 Output directory: {os.path.dirname(result)}")
        else:
            print("❌ Conversion failed!")
    else:
        print(f"❌ File {pptx_file} does not exist!")